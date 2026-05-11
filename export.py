"""
Excel(xlsx) / PowerPoint(pptx) export 모듈.
세션 데이터(분석+발언+다이제스트)를 다운로드 가능한 파일로 변환.
"""

import io
import logging
from typing import Dict, Any, List
from datetime import datetime

log = logging.getLogger(__name__)


AXES = ["데이터", "콘텐츠", "AI Commerce", "UX", "브랜드 메시지 적합도"]


# ============== Excel (xlsx) ==============
def generate_xlsx(session: Dict[str, Any]) -> bytes:
    """세션 전체를 xlsx로 변환."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    bold = Font(bold=True, size=11)
    bold_white = Font(bold=True, size=11, color="FFFFFF")
    title_font = Font(bold=True, size=14)
    header_fill = PatternFill("solid", fgColor="2E5BFF")
    sub_fill = PatternFill("solid", fgColor="F2F2F7")
    thin = Side(border_style="thin", color="DDDDDD")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    wrap = Alignment(wrap_text=True, vertical="top")

    # ----------- Sheet 1: 개요 -----------
    ws = wb.active
    ws.title = "개요"
    ws["A1"] = session.get("title") or session.get("query") or ""
    ws["A1"].font = title_font
    ws.merge_cells("A1:D1")

    meta = [
        ("토론 주제", session.get("query", "")),
        ("모드", "비교 (A/B)" if session.get("mode") == "compare" else "단일"),
        ("URL", session.get("url", "")),
        ("생성일", session.get("created_at", "")),
        ("최대 라운드", session.get("max_rounds", 3)),
        ("페르소나", ", ".join([p.get("name", "") for p in session.get("personas", [])])),
    ]
    if session.get("mode") == "compare":
        meta.append(("좌측(A)", session.get("label_a") or (session.get("side_a") or {}).get("label") or "A"))
        meta.append(("우측(B)", session.get("label_b") or (session.get("side_b") or {}).get("label") or "B"))
        a_url = (session.get("side_a") or {}).get("url", "")
        b_url = (session.get("side_b") or {}).get("url", "")
        if a_url: meta.append(("A URL", a_url))
        if b_url: meta.append(("B URL", b_url))

    for i, (k, v) in enumerate(meta, start=3):
        ws.cell(row=i, column=1, value=k).font = bold
        ws.cell(row=i, column=2, value=str(v))
        ws.cell(row=i, column=1).fill = sub_fill
        ws.merge_cells(start_row=i, start_column=2, end_row=i, end_column=4)

    # AEO 점수
    analysis = session.get("analysis") or {}
    if session.get("mode") == "compare":
        row = len(meta) + 4
        ws.cell(row=row, column=1, value="").value = ""
        for ax_i, ax in enumerate(AXES):
            ws.cell(row=row + 1, column=2 + ax_i, value=ax).font = bold_white
            ws.cell(row=row + 1, column=2 + ax_i).fill = header_fill
            ws.cell(row=row + 1, column=2 + ax_i).alignment = Alignment(horizontal="center")
        ws.cell(row=row + 1, column=1, value="축별 점수").font = bold
        ws.cell(row=row + 1, column=1).fill = sub_fill
        for side_idx, side_key in enumerate(["side_a", "side_b"]):
            side = session.get(side_key) or {}
            an = side.get("analysis") or {}
            label = side.get("label") or ("A" if side_key == "side_a" else "B")
            ws.cell(row=row + 2 + side_idx, column=1, value=label).font = bold
            for ax_i, ax in enumerate(AXES):
                d = (an.get("by_dimension") or {}).get(ax, {})
                ws.cell(row=row + 2 + side_idx, column=2 + ax_i, value=d.get("score", 0)).alignment = Alignment(horizontal="center")
    else:
        row = len(meta) + 4
        ws.cell(row=row, column=1, value="AEO 점수").font = bold
        ws.cell(row=row, column=2, value=f"{analysis.get('aeo_score', 0)}/100")
        ws.cell(row=row + 1, column=1, value="요약").font = bold
        ws.cell(row=row + 1, column=2, value=analysis.get("summary", ""))
        ws.cell(row=row + 2, column=1, value="소비자 인식").font = bold
        ws.cell(row=row + 2, column=2, value=analysis.get("consumer_perception", ""))
        ws.cell(row=row + 1, column=2).alignment = wrap
        ws.cell(row=row + 2, column=2).alignment = wrap
        ws.merge_cells(start_row=row + 1, start_column=2, end_row=row + 1, end_column=4)
        ws.merge_cells(start_row=row + 2, start_column=2, end_row=row + 2, end_column=4)

    ws.column_dimensions["A"].width = 18
    for col in ("B", "C", "D"):
        ws.column_dimensions[col].width = 28

    # ----------- Sheet 2: 축별 분석 -----------
    def _write_dim_sheet(name: str, an: Dict[str, Any]):
        s = wb.create_sheet(name)
        s["A1"] = name
        s["A1"].font = title_font
        s.merge_cells("A1:E1")
        headers = ["축", "점수", "발견(Findings)", "결손(Gaps)", "액션(Actions)"]
        for i, h in enumerate(headers, start=1):
            cell = s.cell(row=3, column=i, value=h)
            cell.font = bold_white
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center", vertical="center")
        for ri, ax in enumerate(AXES, start=4):
            d = (an.get("by_dimension") or {}).get(ax, {})
            s.cell(row=ri, column=1, value=ax).font = bold
            s.cell(row=ri, column=2, value=d.get("score", 0)).alignment = Alignment(horizontal="center")
            s.cell(row=ri, column=3, value="\n".join([f"• {x}" for x in d.get("findings", [])]))
            s.cell(row=ri, column=4, value="\n".join([f"• {x}" for x in d.get("gaps", [])]))
            s.cell(row=ri, column=5, value="\n".join([f"• {x}" for x in d.get("actions", [])]))
            for c in range(1, 6):
                s.cell(row=ri, column=c).alignment = wrap
                s.cell(row=ri, column=c).border = border
            s.row_dimensions[ri].height = 80
        s.column_dimensions["A"].width = 14
        s.column_dimensions["B"].width = 8
        for c in ("C", "D", "E"):
            s.column_dimensions[c].width = 36

    if session.get("mode") == "compare":
        a_an = (session.get("side_a") or {}).get("analysis") or {}
        b_an = (session.get("side_b") or {}).get("analysis") or {}
        a_label = (session.get("side_a") or {}).get("label") or session.get("label_a") or "A"
        b_label = (session.get("side_b") or {}).get("label") or session.get("label_b") or "B"
        _write_dim_sheet(f"축별 분석 - {a_label}"[:31], a_an)
        _write_dim_sheet(f"축별 분석 - {b_label}"[:31], b_an)
    else:
        _write_dim_sheet("축별 분석", analysis)

    # ----------- Sheet 3: 토론 발언 -----------
    s = wb.create_sheet("토론 발언")
    headers = ["라운드", "페르소나", "입장", "대상", "축", "근거", "주장", "액션", "종합"]
    for i, h in enumerate(headers, start=1):
        cell = s.cell(row=1, column=i, value=h)
        cell.font = bold_white
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
    r = 2
    for t in session.get("turns") or []:
        dims = t.get("dimensions") or []
        if not dims:
            s.cell(row=r, column=1, value=t.get("round", ""))
            s.cell(row=r, column=2, value=t.get("persona", ""))
            s.cell(row=r, column=3, value=t.get("stance", ""))
            s.cell(row=r, column=4, value=t.get("target", ""))
            s.cell(row=r, column=9, value=t.get("synthesis", ""))
            r += 1
            continue
        for d in dims:
            s.cell(row=r, column=1, value=t.get("round", ""))
            s.cell(row=r, column=2, value=t.get("persona", ""))
            s.cell(row=r, column=3, value=t.get("stance", ""))
            s.cell(row=r, column=4, value=t.get("target", ""))
            s.cell(row=r, column=5, value=d.get("axis", ""))
            evidence_txt = "\n".join([f"[{e.get('source','')}] {e.get('quote','')}" for e in d.get("evidence", [])])
            s.cell(row=r, column=6, value=evidence_txt)
            s.cell(row=r, column=7, value=d.get("argument", ""))
            s.cell(row=r, column=8, value=d.get("action", ""))
            s.cell(row=r, column=9, value=t.get("synthesis", ""))
            for c in range(1, 10):
                s.cell(row=r, column=c).alignment = wrap
                s.cell(row=r, column=c).border = border
            r += 1
    widths = [8, 16, 10, 14, 12, 36, 50, 36, 30]
    for i, w in enumerate(widths, start=1):
        s.column_dimensions[get_column_letter(i)].width = w
    s.row_dimensions[1].height = 24

    # ----------- Sheet 4: 다이제스트 -----------
    digest_row = session.get("digest") or {}
    digest = digest_row.get("digest") or {}
    s = wb.create_sheet("다이제스트")
    s["A1"] = "라이브 다이제스트"
    s["A1"].font = title_font
    s.merge_cells("A1:D1")
    s["A3"] = "헤드라인"
    s["A3"].font = bold
    s["B3"] = digest.get("headline", "")
    s.merge_cells("B3:D3")

    s["A5"] = "축별 합의·충돌·액션"
    s["A5"].font = bold
    headers = ["축", "합의(Consensus)", "충돌(Conflicts)", "액션(Actions)"]
    for i, h in enumerate(headers, start=1):
        cell = s.cell(row=6, column=i, value=h)
        cell.font = bold_white
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
    for ri, ax in enumerate(AXES, start=7):
        d = (digest.get("by_dimension") or {}).get(ax, {})
        s.cell(row=ri, column=1, value=ax).font = bold
        s.cell(row=ri, column=2, value="\n".join([f"• {x}" for x in d.get("consensus", [])]))
        s.cell(row=ri, column=3, value="\n".join([f"• {x}" for x in d.get("conflicts", [])]))
        s.cell(row=ri, column=4, value="\n".join([f"• {x}" for x in d.get("actions", [])]))
        for c in range(1, 5):
            s.cell(row=ri, column=c).alignment = wrap
            s.cell(row=ri, column=c).border = border
        s.row_dimensions[ri].height = 70

    # 통합 인사이트 + 다음 질문
    base_row = 12
    if digest.get("top_insights"):
        s.cell(row=base_row, column=1, value="통합 인사이트").font = bold
        s.cell(row=base_row, column=2, value="\n".join([f"• {x}" for x in digest["top_insights"]]))
        s.merge_cells(start_row=base_row, start_column=2, end_row=base_row, end_column=4)
        s.cell(row=base_row, column=2).alignment = wrap
        s.row_dimensions[base_row].height = 80
        base_row += 1
    if digest.get("next_questions"):
        s.cell(row=base_row, column=1, value="다음 라운드 질문").font = bold
        s.cell(row=base_row, column=2, value="\n".join([f"• {x}" for x in digest["next_questions"]]))
        s.merge_cells(start_row=base_row, start_column=2, end_row=base_row, end_column=4)
        s.cell(row=base_row, column=2).alignment = wrap
        s.row_dimensions[base_row].height = 70

    s.column_dimensions["A"].width = 18
    for c in ("B", "C", "D"):
        s.column_dimensions[c].width = 32

    # 저장
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ============== PowerPoint (pptx) ==============
def generate_pptx(session: Dict[str, Any]) -> bytes:
    """세션 전체를 pptx로 변환."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLANK = prs.slide_layouts[6]
    NAVY = RGBColor(0x1A, 0x1A, 0x2E)
    BLUE = RGBColor(0x00, 0x7A, 0xFF)
    GRAY_DARK = RGBColor(0x33, 0x33, 0x33)
    GRAY_MID = RGBColor(0x66, 0x66, 0x77)
    GRAY_LIGHT = RGBColor(0xE8, 0xE8, 0xEE)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN = RGBColor(0x30, 0xD1, 0x58)
    ORANGE = RGBColor(0xFF, 0x95, 0x00)
    RED = RGBColor(0xFF, 0x3B, 0x30)

    def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                    color=GRAY_DARK, align=PP_ALIGN.LEFT, font="맑은 고딕"):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text or ""
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tb

    def add_rect(slide, left, top, width, height, fill=GRAY_LIGHT, line=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(left), Inches(top), Inches(width), Inches(height))
        shp.adjustments[0] = 0.08
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(0.75)
        shp.shadow.inherit = False
        return shp

    def score_color(score):
        if score >= 70: return GREEN
        if score >= 40: return ORANGE
        return RED

    # ----------- 슬라이드 1: 표지 -----------
    slide = prs.slides.add_slide(BLANK)
    # 배경 색
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    add_textbox(slide, 0.6, 0.6, 12, 0.5, "AEO Lab", size=14, color=BLUE, bold=True)
    add_textbox(slide, 0.6, 2.3, 12, 1.5, session.get("title") or session.get("query") or "AI 토론 분석 리포트",
                size=40, bold=True, color=WHITE)
    add_textbox(slide, 0.6, 4.0, 12, 1.0, session.get("query", "")[:120],
                size=18, color=RGBColor(0xCC, 0xCC, 0xDD))

    bottom_y = 6.4
    info = []
    if session.get("mode") == "compare":
        a_label = session.get("label_a") or (session.get("side_a") or {}).get("label", "A")
        b_label = session.get("label_b") or (session.get("side_b") or {}).get("label", "B")
        info.append(f"🆚 비교 모드: {a_label} vs {b_label}")
    else:
        info.append("📄 단일 분석 모드")
    if session.get("url"):
        info.append(f"🔗 {session['url'][:80]}")
    info.append(f"📅 {session.get('created_at','')[:10]}")
    info.append(f"👥 {len(session.get('personas') or [])}명 페르소나")
    add_textbox(slide, 0.6, bottom_y, 12, 0.8, "   ·   ".join(info), size=12,
                color=RGBColor(0xCC, 0xCC, 0xDD))

    # ----------- 슬라이드 2: 분석 요약 -----------
    analysis = session.get("analysis") or {}
    slide = prs.slides.add_slide(BLANK)
    add_textbox(slide, 0.6, 0.4, 12, 0.6, "분석 요약", size=28, bold=True, color=NAVY)

    if session.get("mode") == "compare":
        a = session.get("side_a") or {}
        b = session.get("side_b") or {}
        a_an = a.get("analysis") or {}
        b_an = b.get("analysis") or {}
        a_label = a.get("label") or session.get("label_a") or "A"
        b_label = b.get("label") or session.get("label_b") or "B"

        for ci, (label, an) in enumerate([(a_label, a_an), (b_label, b_an)]):
            x = 0.6 + ci * 6.3
            add_rect(slide, x, 1.3, 6.0, 5.7, fill=GRAY_LIGHT)
            add_textbox(slide, x + 0.3, 1.5, 5.4, 0.6, label, size=22, bold=True, color=NAVY)
            add_textbox(slide, x + 0.3, 2.2, 5.4, 0.5, f"AEO {an.get('aeo_score', 0)}/100",
                        size=18, bold=True, color=score_color(an.get('aeo_score', 0)))
            add_textbox(slide, x + 0.3, 2.9, 5.4, 1.2, an.get("summary", ""), size=12, color=GRAY_DARK)
            # 축별 미니 점수
            for ax_i, ax in enumerate(AXES):
                d = (an.get("by_dimension") or {}).get(ax, {})
                cy = 4.3 + ax_i * 0.55
                add_textbox(slide, x + 0.3, cy, 2.5, 0.4, ax, size=12, bold=True, color=GRAY_DARK)
                add_textbox(slide, x + 3.0, cy, 2.0, 0.4, f"{d.get('score',0)}",
                            size=14, bold=True, color=score_color(d.get('score', 0)), align=PP_ALIGN.RIGHT)
    else:
        add_textbox(slide, 0.6, 1.3, 4, 0.6,
                    f"AEO 점수  {analysis.get('aeo_score',0)}/100",
                    size=22, bold=True, color=score_color(analysis.get('aeo_score', 0)))
        add_textbox(slide, 0.6, 2.1, 12, 1.0, analysis.get("summary", ""), size=14, color=GRAY_DARK)

        # 축별 카드 5개
        for ax_i, ax in enumerate(AXES):
            x = 0.6 + (ax_i % 2) * 6.3
            y = 3.3 + (ax_i // 2) * 2.0
            d = (analysis.get("by_dimension") or {}).get(ax, {})
            add_rect(slide, x, y, 6.0, 1.85, fill=GRAY_LIGHT)
            add_textbox(slide, x + 0.25, y + 0.15, 4, 0.5, ax, size=16, bold=True, color=NAVY)
            add_textbox(slide, x + 4.5, y + 0.15, 1.3, 0.5,
                        f"{d.get('score',0)}", size=22, bold=True,
                        color=score_color(d.get('score', 0)), align=PP_ALIGN.RIGHT)
            actions_text = " · ".join(d.get("actions", [])[:2])[:130]
            add_textbox(slide, x + 0.25, y + 0.85, 5.5, 1.0, actions_text, size=11, color=GRAY_MID)

    # ----------- 슬라이드 3+: 발언 -----------
    turns = session.get("turns") or []
    if turns:
        # 페르소나별로 묶거나, 그냥 라운드별로
        for t in turns:
            slide = prs.slides.add_slide(BLANK)
            persona = t.get("persona", "?")
            emoji = t.get("emoji", "💬")
            stance = t.get("stance", "")
            target = t.get("target", "")
            round_no = t.get("round", "")
            synthesis = t.get("synthesis", "")

            # 헤더
            add_textbox(slide, 0.6, 0.4, 8, 0.6, f"{emoji}  {persona}",
                        size=22, bold=True, color=NAVY)
            tag = f"[{stance}]"
            if target:
                tag += f"  → {target}"
            add_textbox(slide, 8.6, 0.55, 4, 0.5, f"Round {round_no}    {tag}",
                        size=12, color=GRAY_MID, align=PP_ALIGN.RIGHT)

            if synthesis:
                add_rect(slide, 0.6, 1.1, 12.1, 0.7, fill=RGBColor(0xEB, 0xF2, 0xFF))
                add_textbox(slide, 0.85, 1.25, 11.7, 0.5, f"\u201c{synthesis}\u201d",
                            size=14, bold=True, color=NAVY)

            # 차원별 카드
            dims = t.get("dimensions") or []
            base_y = 2.05
            for di, d in enumerate(dims[:2]):  # 한 슬라이드 최대 2 dim
                y = base_y + di * 2.55
                add_rect(slide, 0.6, y, 12.1, 2.45, fill=GRAY_LIGHT)
                add_textbox(slide, 0.85, y + 0.15, 8, 0.5,
                            f"📐 {d.get('axis','')}", size=14, bold=True, color=BLUE)
                # 근거 (오른쪽)
                ev_txt = "\n".join([f"• [{e.get('source','')}] {e.get('quote','')[:80]}" for e in d.get("evidence", [])[:3]])
                add_textbox(slide, 0.85, y + 0.7, 5.5, 1.6, ev_txt or "(근거 없음)",
                            size=10, color=GRAY_DARK)
                # 주장
                add_textbox(slide, 6.6, y + 0.7, 6.0, 1.2,
                            d.get("argument", "")[:280], size=11, color=GRAY_DARK)
                # 액션 (작은 박스)
                add_rect(slide, 6.6, y + 1.95, 6.0, 0.4, fill=RGBColor(0xE0, 0xFA, 0xE6))
                add_textbox(slide, 6.75, y + 2.0, 5.7, 0.35,
                            f"✓ {d.get('action','')[:130]}", size=10, color=GRAY_DARK)

    # ----------- 마지막: 다이제스트 -----------
    digest = (session.get("digest") or {}).get("digest") or {}
    if digest:
        slide = prs.slides.add_slide(BLANK)
        add_textbox(slide, 0.6, 0.4, 12, 0.6, "라이브 다이제스트", size=28, bold=True, color=NAVY)
        if digest.get("headline"):
            add_rect(slide, 0.6, 1.2, 12.1, 0.7, fill=RGBColor(0xEB, 0xF2, 0xFF))
            add_textbox(slide, 0.85, 1.35, 11.7, 0.5,
                        f"\u201c{digest['headline']}\u201d", size=15, bold=True, color=NAVY)

        # 5개 축
        for ax_i, ax in enumerate(AXES):
            x = 0.6 + (ax_i % 2) * 6.3
            y = 2.2 + (ax_i // 2) * 2.55
            d = (digest.get("by_dimension") or {}).get(ax, {})
            add_rect(slide, x, y, 6.0, 2.45, fill=GRAY_LIGHT)
            add_textbox(slide, x + 0.2, y + 0.15, 5, 0.5, ax, size=14, bold=True, color=NAVY)
            sub_y = y + 0.7
            for label, color, items in [
                ("합의", GREEN, d.get("consensus", [])),
                ("충돌", ORANGE, d.get("conflicts", [])),
                ("액션", BLUE, d.get("actions", [])),
            ]:
                if not items:
                    continue
                add_textbox(slide, x + 0.2, sub_y, 1, 0.3, label, size=10, bold=True, color=color)
                txt = " · ".join(items[:2])[:110]
                add_textbox(slide, x + 1.0, sub_y, 4.8, 0.5, txt, size=9, color=GRAY_DARK)
                sub_y += 0.55

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
